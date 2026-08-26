"""AUCE 2026 senior-project defense PowerPoint — cinematic, click-driven, academic."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(r"c:\Users\mohamad\.vscode\Senior_Project_AUCE_2026\AI_Network_Analyzer")
SHOTS = ROOT / "research_paper_assets" / "user_guide_shots"
EVAL = ROOT / "research_paper_assets" / "evaluation"
LOGO = ROOT / "dashboard" / "assets" / "app_logo.png"
MEDIA = ROOT / "research_paper_assets" / "ppt_media"
OUT_DESK = Path(r"C:\Users\mohamad\OneDrive\Desktop\AI_Network_Analyzer_Senior_Project_Defense.pptx")
OUT_PROJ = ROOT / "research_paper_assets" / "AI_Network_Analyzer_Senior_Project_Defense.pptx"

NAVY = RGBColor(0x0C, 0x1A, 0x2E)
NAVY2 = RGBColor(0x10, 0x24, 0x42)
CARD = RGBColor(0x14, 0x2C, 0x4C)
GOLD = RGBColor(0xF0, 0xC1, 0x4B)
GOLD2 = RGBColor(0xD4, 0xA0, 0x2A)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x9B, 0xB0, 0xC9)
CYAN = RGBColor(0x5B, 0xB8, 0xFF)
GREEN = RGBColor(0x3D, 0xDC, 0x97)
RED = RGBColor(0xFF, 0x6B, 0x7A)

W, H = 13.333333, 7.5
NSMAP_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
NSMAP_P = "http://schemas.openxmlformats.org/presentationml/2006/main"


def crop169(src: Path, dst: Path) -> Path:
    dst.parent.mkdir(parents=True, exist_ok=True)
    if not src.exists() and dst.exists():
        return dst
    im = Image.open(src).convert("RGB")
    w, h = im.size
    th = int(w * 9 / 16)
    if h > th:
        im = im.crop((0, 0, w, th))
    im.save(dst, quality=92)
    return dst


def make_flow_art() -> Path:
    MEDIA.mkdir(parents=True, exist_ok=True)
    path = MEDIA / "flow_art.png"
    img = Image.new("RGB", (1600, 420), (12, 26, 46))
    d = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("calibri.ttf", 28)
        small = ImageFont.truetype("calibri.ttf", 18)
    except Exception:
        font = ImageFont.load_default()
        small = font
    steps = [
        ("1", "Capture", "Live packets"),
        ("2", "Features", "39 flow stats"),
        ("3", "Hybrid AI", "5 models"),
        ("4", "Decide", "Score + label"),
        ("5", "SOC", "Alert / incident"),
        ("6", "Act", "Block · XAI"),
    ]
    for i, (n, title, sub) in enumerate(steps):
        x = 40 + i * 260
        d.rounded_rectangle((x, 70, x + 220, 320), 22, fill=(20, 44, 76), outline=(240, 193, 75), width=2)
        d.ellipse((x + 88, 95, x + 132, 139), fill=(240, 193, 75))
        d.text((x + 104, 102), n, fill=(12, 26, 46), font=font)
        d.text((x + 28, 165), title, fill=(248, 250, 252), font=font)
        d.text((x + 28, 215), sub, fill=(155, 176, 201), font=small)
        if i < 5:
            d.polygon([(x + 228, 190), (x + 248, 205), (x + 228, 220)], fill=(240, 193, 75))
    img.save(path)
    return path


def make_ai_art() -> Path:
    path = MEDIA / "hybrid_ai.png"
    img = Image.new("RGB", (1600, 520), (12, 26, 46))
    d = ImageDraw.Draw(img)
    try:
        font = ImageFont.truetype("calibri.ttf", 26)
        small = ImageFont.truetype("calibri.ttf", 18)
        big = ImageFont.truetype("calibrib.ttf", 32)
    except Exception:
        font = ImageFont.load_default()
        small = font
        big = font
    models = [
        ("Random Forest", "Supervised", "Tabular attacks"),
        ("XGBoost", "Supervised", "High precision"),
        ("Isolation Forest", "Anomaly", "Unknown traffic"),
        ("Autoencoder", "Deep", "Reconstruction"),
        ("LSTM", "Sequence", "Time patterns"),
    ]
    for i, (title, kind, role) in enumerate(models):
        x = 30 + i * 314
        d.rounded_rectangle((x, 40, x + 290, 280), 20, fill=(20, 44, 76), outline=(91, 184, 255), width=2)
        d.text((x + 18, 70), title, fill=(240, 193, 75), font=font)
        d.text((x + 18, 120), kind, fill=(91, 184, 255), font=small)
        d.text((x + 18, 170), role, fill=(248, 250, 252), font=small)
    d.rounded_rectangle((220, 330, 1380, 490), 20, fill=(16, 36, 66), outline=(240, 193, 75), width=2)
    d.text((250, 365), "Hybrid decision  →  threat score  →  severity  →  SOAR playbook", fill=(248, 250, 252), font=big)
    d.text((250, 420), "Online SGD learns from confirmed attacks without replacing core models", fill=(155, 176, 201), font=small)
    img.save(path)
    return path


class Deck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self.blank = self.prs.slide_layouts[6]
        self.slides = []
        self.notes = []

    def slide(self, section: str, notes: str = ""):
        s = self.prs.slides.add_slide(self.blank)
        bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY
        bg.line.fill.background()
        bg.name = "chrome_bg"
        rail = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), Inches(H))
        rail.fill.solid()
        rail.fill.fore_color.rgb = GOLD
        rail.line.fill.background()
        rail.name = "chrome_rail"
        top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(0.04))
        top.fill.solid()
        top.fill.fore_color.rgb = GOLD
        top.line.fill.background()
        top.name = "chrome_top"
        self.slides.append((s, section))
        self.notes.append(notes)
        return s

    def footer(self, s, idx: int, total: int, section: str):
        line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(7.18), Inches(12.5), Inches(0.015))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x2A, 0x45, 0x6B)
        line.line.fill.background()
        line.name = "chrome_line"
        self.text(s, 0.4, 7.22, 9.2, 0.25, f"AI Network Analyzer  ·  {section}  ·  AUCE 2026  ·  v1.0.0", 11, color=MUTED, name="chrome_f")
        self.text(s, 11.6, 7.22, 1.3, 0.25, f"{idx:02d}  /  {total:02d}", 11, color=GOLD, align=PP_ALIGN.RIGHT, name="chrome_n")

    def text(self, s, l, t, w, h, msg, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, name=None, font="Calibri"):
        box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = msg
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
        rPr = run._r.get_or_add_rPr()
        rFonts = rPr.find(qn("a:latin"))
        if rFonts is None:
            rFonts = etree.SubElement(rPr, qn("a:latin"))
        rFonts.set("typeface", font)
        if name:
            box.name = name
        return box

    def multilines(self, s, l, t, w, h, lines, size=16, color=WHITE, name=None, spacing=8):
        box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.space_after = Pt(spacing)
            run = p.add_run()
            run.text = line
            run.font.size = Pt(size)
            run.font.color.rgb = color
            run.font.name = "Calibri"
        if name:
            box.name = name
        return box

    def rect(self, s, l, t, w, h, fill=CARD, line=None, name=None, round=True):
        shp = s.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,
            Inches(l), Inches(t), Inches(w), Inches(h),
        )
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if line:
            shp.line.color.rgb = line
            shp.line.width = Pt(1.25)
        else:
            shp.line.fill.background()
        if round:
            try:
                shp.adjustments[0] = 0.08
            except Exception:
                pass
        if name:
            shp.name = name
        return shp

    def accent_card(self, s, l, t, w, h, name=None):
        card = self.rect(s, l, t, w, h, CARD, name=name)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(0.07), Inches(h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = GOLD
        bar.line.fill.background()
        bar.name = "with_" + (name or "card") + "_bar"
        return card

    def picture(self, s, path, l, t, w, name=None):
        pic = s.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))
        if name:
            pic.name = name
        return pic

    def kicker(self, s, text, l=0.45, t=0.22):
        return self.text(s, l, t, 8, 0.32, text.upper(), 12, True, GOLD, name="auto_kicker")

    def h1(self, s, text, l=0.45, t=0.48, w=12.4):
        return self.text(s, l, t, w, 0.7, text, 32, True, WHITE, name="auto_title")


def add_fade_transition(slide):
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    trans = etree.Element(qn("p:transition"))
    trans.set("advClick", "1")
    trans.set("spd", "med")
    trans.set("{%s}dur" % NSMAP_P14, "500")
    push = etree.SubElement(trans, qn("p:push"))
    push.set("dir", "r")
    sld.append(trans)


def apply_notes(prs, notes):
    for slide, note in zip(prs.slides, notes):
        if not note:
            continue
        slide.notes_slide.notes_text_frame.text = note


def build():
    MEDIA.mkdir(parents=True, exist_ok=True)
    flow = make_flow_art()
    hybrid = make_ai_art()
    shots = {}
    for key, src in {
        "login": "00_login.png",
        "home": "00_home.png",
        "live": "01_live.png",
        "detect": "02_detection.png",
        "sim": "03_simulation.png",
        "ti": "04_ti.png",
        "mitre": "10_soc_tab0_MITRE_ATT_CK.png",
        "soar": "10_soc_tab1_SOAR_Playbooks.png",
        "copilot": "15_copilot.png",
        "clear": "09_settings_tab4_Clear_Data.png",
        "models": "07_models.png",
        "resp": "14_response.png",
        "alerts": "05_alerts.png",
        "blocked": "06_blocked.png",
        "reports": "08_reports.png",
        "settings": "09_settings.png",
        "incidents": "11_incidents.png",
        "assets": "12_assets.png",
        "hunting": "13_hunting.png",
    }.items():
        shots[key] = crop169(SHOTS / src, MEDIA / f"{key}.jpg")

    d = Deck()
    S = []

    def add(section, notes=""):
        s = d.slide(section, notes)
        S.append(s)
        return s

    # 1 Title
    s = add("Title", "Pause 3 seconds. State your name, AUCE, and the project as an NDR platform — not only a classifier.")
    if LOGO.exists():
        d.picture(s, LOGO, 0.5, 1.15, 1.35, "auto_logo")
    d.text(s, 2.05, 1.2, 10.5, 0.35, "AMERICAN UNIVERSITY OF CULTURE AND EDUCATION  ·  SENIOR PROJECT 2026", 13, True, GOLD, name="auto_uni")
    d.text(s, 2.05, 1.65, 10.7, 1.35, "AI-Powered Network Traffic Analyser\n& Anomaly Detector", 36, True, WHITE, name="auto_title")
    d.text(s, 2.05, 3.15, 10.5, 0.4, "NDR Platform v1.0.0  —  Detect  ·  Understand  ·  Respond  ·  Explain", 18, False, CYAN, name="auto_sub")
    d.rect(s, 0.5, 4.15, 12.3, 0.02, GOLD, round=False, name="auto_rule")
    d.accent_card(s, 0.5, 4.5, 3.9, 1.7, "click_meta1")
    d.text(s, 0.75, 4.68, 3.5, 0.3, "STUDENT", 11, True, GOLD, name="with_m1")
    d.text(s, 0.75, 5.05, 3.5, 0.9, "Mahmoud Talal Kanaan\nComputer Science", 16, False, WHITE, name="with_m1b")
    d.accent_card(s, 4.7, 4.5, 3.9, 1.7, "click_meta2")
    d.text(s, 4.95, 4.68, 3.5, 0.3, "SUPERVISOR", 11, True, GOLD, name="with_m2")
    d.text(s, 4.95, 5.05, 3.5, 0.9, "Dr. Hassan Noureddine\nAUCE", 16, False, WHITE, name="with_m2b")
    d.accent_card(s, 8.9, 4.5, 3.9, 1.7, "click_meta3")
    d.text(s, 9.15, 4.68, 3.5, 0.3, "DEFENSE", 11, True, GOLD, name="with_m3")
    d.text(s, 9.15, 5.05, 3.5, 0.9, "Implementation + Evaluation\nLocal SOC / NDR prototype", 16, False, WHITE, name="with_m3b")

    # 2 Agenda
    s = add("Agenda", "Same order as the research paper. Chapter 1 to conclusion.")
    d.kicker(s, "Roadmap  ·  same chapters as the paper")
    d.h1(s, "The defense follows the thesis.")
    items = [
        ("01", "Ch. 1 Introduction", "Problem, objectives O1–O6, proposed solution"),
        ("02", "Ch. 2 Literature", "Twenty platforms and the remaining gap"),
        ("03", "Ch. 3 Architecture", "Pipeline, hybrid AI, MITRE, SOAR"),
        ("04", "Ch. 4 Implementation", "Stack and all 16 SOC pages"),
        ("05", "Ch. 5 Evaluation", "12,000-row holdout — Table 9 / Table 11"),
        ("06", "Ch. 6–8 Close", "Deploy, limits, demo, conclusion"),
    ]
    for i, (n, title, sub) in enumerate(items):
        col, row = i % 3, i // 3
        x, y = 0.45 + col * 4.2, 1.45 + row * 2.55
        d.accent_card(s, x, y, 3.95, 2.25, f"click_ag{i}")
        d.text(s, x + 0.25, y + 0.25, 3.4, 0.45, n, 28, True, GOLD, name=f"with_ag{i}n")
        d.text(s, x + 0.25, y + 0.85, 3.4, 0.5, title, 22, True, WHITE, name=f"with_ag{i}t")
        d.text(s, x + 0.25, y + 1.4, 3.4, 0.6, sub, 14, False, MUTED, name=f"with_ag{i}s")

    # 3 Problem
    s = add("Problem", "Do not list tools. State the pain: signatures miss novelty; NDR is expensive; labs have no SOC loop.")
    d.kicker(s, "Ch. 1  ·  Introduction")
    d.h1(s, "Attacks are live. Most student labs still watch logs.")
    pains = [
        ("Signature IDS", "Snort-style rules fail on unseen variants and encrypted or novel campaigns."),
        ("AI without SOC", "A classifier that cannot alert, ticket, hunt, or block is not operations."),
        ("Enterprise NDR cost", "Commercial platforms are out of reach for labs, SMEs, and teaching."),
        ("Broken loop", "Capture, detect, notify, and respond usually live in separate products."),
    ]
    for i, (t, b) in enumerate(pains):
        x, y = 0.45 + (i % 2) * 6.4, 1.45 + (i // 2) * 2.5
        d.accent_card(s, x, y, 6.15, 2.25, f"click_p{i}")
        d.text(s, x + 0.3, y + 0.3, 5.6, 0.5, t, 22, True, GOLD, name=f"with_p{i}t")
        d.text(s, x + 0.3, y + 0.95, 5.6, 1.0, b, 16, False, WHITE, name=f"with_p{i}b")

    # 4 Motivation
    s = add("Motivation", "Why AUCE should care: a teaching-grade NDR that is implemented, not only designed.")
    d.kicker(s, "Ch. 1  ·  Why this project")
    d.h1(s, "Build a SOC you can run on a laptop — and defend with evidence.")
    points = [
        ("Academic need", "Supervisor guidance: move from design-only writing to a running system plus evaluation."),
        ("Operational need", "Analysts need one glass: traffic, detections, incidents, hunting, and response."),
        ("Pedagogical need", "Students must see the full kill-chain loop: inject → detect → alert → block → explain."),
        ("Local control", "On-prem FastAPI + Streamlit + SQLite (Postgres optional). No cloud lock-in for the prototype."),
    ]
    for i, (t, b) in enumerate(points):
        y = 1.4 + i * 1.3
        d.accent_card(s, 0.45, y, 12.4, 1.15, f"click_mot{i}")
        d.text(s, 0.75, y + 0.12, 11.8, 0.35, t, 18, True, GOLD, name=f"with_mot{i}t")
        d.text(s, 0.75, y + 0.52, 11.8, 0.5, b, 15, False, WHITE, name=f"with_mot{i}b")

    # Problem vs solution (paper 1.2 / 1.3)
    s = add("Problem-solution", "Read left then right. Each pain has a matching capability.")
    d.kicker(s, "Ch. 1  ·  Problem and solution")
    d.h1(s, "Four pains. Four answers in the product.")
    pairs = [
        ("Signature IDS miss novel traffic", "Hybrid models score known and unknown flows"),
        ("AI without a SOC workflow", "16 pages: alerts, incidents, hunting, reports"),
        ("Enterprise NDR is too costly", "Runs locally: FastAPI + Streamlit + SQLite (Postgres optional)"),
        ("Detect with no response", "SOAR, block, Telegram, XAI, Copilot, approvals"),
    ]
    d.rect(s, 0.45, 1.35, 6.1, 0.5, RGBColor(0x3A, 0x1B, 0x24), name="auto_ph")
    d.text(s, 0.6, 1.42, 5.8, 0.4, "PROBLEM", 14, True, GOLD, name="with_ph")
    d.rect(s, 6.75, 1.35, 6.1, 0.5, RGBColor(0x1A, 0x3A, 0x28), name="auto_sh")
    d.text(s, 6.9, 1.42, 5.8, 0.4, "SOLUTION IN THIS PROJECT", 14, True, GOLD, name="with_sh")
    for i, (p, sol) in enumerate(pairs):
        y = 2.0 + i * 1.15
        d.accent_card(s, 0.45, y, 6.1, 1.0, f"click_ps{i}")
        d.text(s, 0.7, y + 0.28, 5.6, 0.5, p, 16, False, WHITE, name=f"with_ps{i}p")
        d.rect(s, 6.75, y, 6.1, 1.0, CARD, GREEN, name=f"click_ss{i}")
        d.text(s, 7.0, y + 0.28, 5.6, 0.5, sol, 16, False, WHITE, name=f"with_ss{i}s")

    # Objectives — same O1–O6 as paper Table 2
    s = add("Objectives", "Read the same six objectives as Table 2 in the paper.")
    d.kicker(s, "Ch. 1  ·  Research objectives")
    d.h1(s, "Six objectives. Same O1–O6 as the paper.")
    objs = [
        ("O1", "Monitor network traffic and extract useful security features (39 CICIDS-style flow features)."),
        ("O2", "Detect known attacks using supervised models: Random Forest and XGBoost."),
        ("O3", "Detect unknown anomalies using Isolation Forest and Autoencoder."),
        ("O4", "Use LSTM to analyze sequential behavior and support attack prediction."),
        ("O5", "Improve transparency using stored per-prediction XAI (this-flow evidence + tree importances)."),
        ("O6", "Support incident response: IP blocking, Telegram, webhooks, and reports (no email)."),
    ]
    for i, (k, t) in enumerate(objs):
        y = 1.28 + i * 0.88
        d.rect(s, 0.45, y, 1.15, 0.75, GOLD, name=f"click_o{i}")
        d.text(s, 0.45, y + 0.16, 1.15, 0.45, k, 20, True, NAVY, align=PP_ALIGN.CENTER, name=f"with_o{i}k")
        d.accent_card(s, 1.8, y, 11.05, 0.75, f"with_o{i}c")
        d.text(s, 2.05, y + 0.16, 10.6, 0.48, t, 15, False, WHITE, name=f"with_o{i}t")

    # Proposed solution
    s = add("Solution", "One sentence: a local NDR that detects with hybrid AI and operates like a mini SOC.")
    d.kicker(s, "Ch. 1  ·  Proposed solution")
    d.h1(s, "One product. Four verbs.")
    pillars = [
        ("DETECT", "Live capture + hybrid models + threat score + severity."),
        ("UNDERSTAND", "MITRE ATT&CK, threat intel, host risk, stored XAI, Copilot."),
        ("RESPOND", "SOAR playbooks, allow/block, approvals, Telegram, SIEM/Jira hooks."),
        ("PROVE", "Reports, live-versus-simulation evidence, and the 12k holdout."),
    ]
    for i, (t, b) in enumerate(pillars):
        x = 0.45 + i * 3.2
        d.rect(s, x, 1.5, 3.05, 4.7, CARD, GOLD, name=f"click_pil{i}")
        d.text(s, x + 0.2, 1.8, 2.65, 0.7, f"0{i+1}", 28, True, GOLD, name=f"with_pil{i}n")
        d.text(s, x + 0.2, 2.6, 2.65, 0.8, t, 22, True, WHITE, name=f"with_pil{i}t")
        d.text(s, x + 0.2, 3.5, 2.65, 2.2, b, 15, False, MUTED, name=f"with_pil{i}b")

    # Ch.2 Gap
    s = add("Gap", "Position against Snort, Wazuh, and cloud NDR. Your gap is integration + hybrid AI + local SOC.")
    d.kicker(s, "Ch. 2  ·  Literature gap")
    d.h1(s, "Existing tools solve a slice. This project closes the loop.")
    headers = ["Approach", "Strength", "What is still missing"]
    rows = [
        ["Signature IDS", "Fast known threats", "Weak on novel / zero-day patterns"],
        ["SIEM / EDR", "Logs & endpoints", "Not a packet-level NDR with hybrid AI"],
        ["Cloud NDR", "Scale & intel", "Cost, privacy, not a student-lab stack"],
        ["This work", "Local NDR + SOC", "Implemented pipeline + measured models"],
    ]
    for i, h in enumerate(headers):
        d.rect(s, 0.45 + i * 4.2, 1.4, 4.05, 0.55, RGBColor(0x1B, 0x3A, 0x5F), name=f"auto_head{i}")
        d.text(s, 0.55 + i * 4.2, 1.48, 3.85, 0.4, h, 14, True, GOLD, name=f"with_head{i}")
    for r, row in enumerate(rows):
        y = 2.1 + r * 1.1
        fill = RGBColor(0x1A, 0x48, 0x38) if r == 3 else CARD
        d.rect(s, 0.45, y, 12.45, 1.0, fill, name=f"click_gap{r}")
        for c, cell in enumerate(row):
            d.text(s, 0.6 + c * 4.2, y + 0.28, 3.9, 0.5, cell, 15, r == 3, WHITE if r == 3 else MUTED, name=f"with_gap{r}{c}")

    # Why different
    s = add("Difference", "This is the 'why us' slide. Four differences. Do not rush.")
    d.kicker(s, "Ch. 2  ·  How the gap is closed")
    d.h1(s, "One local NDR. Capabilities usually split across several tools.")
    diffs = [
        ("End-to-end loop", "Capture, detect, alert, hunt, and block in one prototype — typically split across IDS, SIEM, and a notebook."),
        ("Hybrid AI", "Supervised + anomaly + deep + sequence + online learning that does not overwrite core models."),
        ("Same path for demos", "Threat simulation injects into the real pipeline. The jury sees true SOC behaviour."),
        ("Operate, then explain", "MITRE, SOAR, per-alert XAI, Copilot, Telegram, and human-approval mode — not accuracy alone."),
    ]
    for i, (t, b) in enumerate(diffs):
        x, y = 0.45 + (i % 2) * 6.4, 1.45 + (i // 2) * 2.5
        d.accent_card(s, x, y, 6.15, 2.25, f"click_df{i}")
        d.text(s, x + 0.3, y + 0.3, 5.6, 0.5, t, 22, True, GOLD, name=f"with_df{i}t")
        d.text(s, x + 0.3, y + 0.95, 5.6, 1.0, b, 16, False, WHITE, name=f"with_df{i}b")

    # 8 Architecture
    s = add("Architecture", "Walk left to right. Click each stage. This is the backbone of the oral defense.")
    d.kicker(s, "Ch. 3  ·  Architecture")
    d.h1(s, "From packet to playbook — on one machine.")
    d.picture(s, flow, 0.45, 1.5, 12.4, "click_flow")
    d.text(s, 0.5, 5.55, 12.3, 1.2,
           "run.bat starts FastAPI (:8000) and Streamlit (:8501). SQLite stores flows, predictions, alerts, XAI, blocks, and incidents. PostgreSQL is optional. A supervisor can auto-recover; run_capture.bat is the 24/7 sensor.",
           15, False, MUTED, name="click_flowcap")

    # 9 Hybrid AI
    s = add("Hybrid AI", "Emphasize hybrid: supervised + anomaly + deep + sequence + online learning that does not overwrite cores.")
    d.kicker(s, "Ch. 3  ·  Hybrid AI engine")
    d.h1(s, "Five brains. One decision.")
    d.picture(s, hybrid, 0.4, 1.35, 12.5, "click_hyb")

    # 10 Tech stack (section 03 — before the platform tour)
    s = add("Stack", "Keep it visual. Python is the spine. FastAPI + Streamlit + SQLite is the body.")
    d.kicker(s, "Ch. 4  ·  Technology stack")
    d.h1(s, "A stack a senior can defend.")
    stack = [
        ("Python 3", "Core language"),
        ("FastAPI", "REST API :8000"),
        ("Streamlit", "SOC UI :8501"),
        ("SQLite", "Local store · Postgres optional"),
        ("Scikit-learn", "RF / Isolation Forest"),
        ("XGBoost", "Gradient boosting"),
        ("PyTorch", "AE / LSTM"),
        ("Scapy", "Live capture"),
        ("Telegram", "Analyst notify"),
        ("MITRE map", "ATT&CK linkage"),
        ("SOAR", "Playbook actions"),
        ("Reports", "CSV / JSON export"),
    ]
    for i, (a, b) in enumerate(stack):
        x, y = 0.45 + (i % 4) * 3.2, 1.4 + (i // 4) * 1.75
        prefix = "click" if i % 4 == 0 else "with"
        d.accent_card(s, x, y, 3.05, 1.55, f"{prefix}_st{i}")
        d.text(s, x + 0.2, y + 0.25, 2.65, 0.5, a, 18, True, WHITE, name=f"with_st{i}a")
        d.text(s, x + 0.2, y + 0.85, 2.65, 0.4, b, 13, False, MUTED, name=f"with_st{i}b")

    def shot_pair(section, kicker, title, notes, left, lcap, right, rcap):
        s = add(section, notes)
        d.kicker(s, kicker)
        d.h1(s, title)
        d.rect(s, 0.4, 1.35, 6.15, 5.15, RGBColor(0x08, 0x12, 0x22), GOLD, name="click_a")
        d.picture(s, shots[left], 0.55, 1.7, 5.85, "with_ap")
        d.text(s, 0.55, 5.95, 5.85, 0.35, lcap, 13, False, MUTED, name="with_ac")
        d.rect(s, 6.75, 1.35, 6.15, 5.15, RGBColor(0x08, 0x12, 0x22), GOLD, name="click_b")
        d.picture(s, shots[right], 6.9, 1.7, 5.85, "with_bp")
        d.text(s, 6.9, 5.95, 5.85, 0.35, rcap, 13, False, MUTED, name="with_bc")
        return s

    def catalog(kicker, title, notes, rows, tag):
        s = add("SOC catalog", notes)
        d.kicker(s, kicker)
        d.h1(s, title)
        for i, (num, name, role) in enumerate(rows):
            y = 1.32 + i * 0.68
            d.rect(s, 0.45, y, 0.9, 0.58, GOLD, name=f"click_{tag}{i}")
            d.text(s, 0.45, y + 0.1, 0.9, 0.4, num, 18, True, NAVY, align=PP_ALIGN.CENTER, name=f"with_{tag}{i}n")
            d.accent_card(s, 1.5, y, 11.35, 0.58, f"with_{tag}{i}c")
            d.text(s, 1.75, y + 0.04, 3.4, 0.48, name, 16, True, WHITE, name=f"with_{tag}{i}t")
            d.text(s, 5.2, y + 0.1, 7.4, 0.42, role, 14, False, MUTED, name=f"with_{tag}{i}r")
        return s

    # 11-12 Full SOC catalog
    catalog(
        "Ch. 4  ·  Dashboard pages 01–08",
        "Every page, with its job.",
        "Name each page and its job. Home through AI Models.",
        [
            ("01", "Home", "SOC pulse: KPIs, charts, shortcuts, recent alerts and blocks"),
            ("02", "Live Monitoring", "Select NIC, start/stop capture, live-evidence snapshot"),
            ("03", "AI Detection", "Hybrid inference, labels, threat score, stored per-flow XAI"),
            ("04", "Threat Simulation", "Inject lab campaigns into the same live detection path"),
            ("05", "Threat Intelligence", "IP reputation, geo, stored intel records"),
            ("06", "Alerts", "Analyst inbox: status, MITRE, true/false-positive notes"),
            ("07", "Blocked IPs", "Manual and automatic block, reason, unblock"),
            ("08", "AI Models", "Accuracy table, retrain from CSV, online learning status"),
        ],
        "c1",
    )
    catalog(
        "Ch. 4  ·  Dashboard pages 09–16",
        "Operations, hunting, and response.",
        "Finish the 16. Tabs matter: SOC Ops, Settings, Reports, Assets, Hunting, Response.",
        [
            ("09", "Reports", "Detection, Incidents, Export, Compliance, Live Evidence"),
            ("10", "Settings", "General, Telegram, AI Config, NDR/Response, Clear Data, Company"),
            ("11", "SOC Ops", "MITRE, SOAR playbooks, online learning, server health, ML registry"),
            ("12", "Incidents", "Group related alerts into one case; owner, status, recommendation"),
            ("13", "Assets", "Inventory, risk, topology, host profile, identity / UEBA"),
            ("14", "Hunting", "Query by IP/attack/protocol; PCAP; DNS/TLS forensics"),
            ("15", "Response", "IOC, allow/block, approvals, webhooks, distributed sensors"),
            ("16", "Copilot", "Natural-language questions, incident summary, per-alert XAI"),
        ],
        "c2",
    )

    shot_pair("Live path", "Ch. 4  ·  Page 01", "Login, then the SOC pulse.",
              "Administrator session. Home is page 01 — the operator glass.",
              "login", "Login  ·  v1.0.0  ·  administrator", "home", "01  Home  ·  KPIs, charts, shortcuts")
    shot_pair("Capture & AI", "Ch. 4  ·  Pages 02–03", "Traffic in. Verdicts out.",
              "Live NIC capture feeds hybrid detection.",
              "live", "02  Live Monitoring  ·  interface and flows", "detect", "03  AI Detection  ·  hybrid inference")
    shot_pair("Simulation", "Ch. 4  ·  Pages 04–05", "Inject a campaign. Enrich the IP.",
              "Simulation uses the real pipeline. TI is reputation plus geo.",
              "sim", "04  Threat Simulation  ·  lab campaigns", "ti", "05  Threat Intelligence  ·  reputation + geo")
    shot_pair("Alerts & blocks", "Ch. 4  ·  Pages 06–07", "The inbox and the block list.",
              "Alerts are the analyst queue. Blocked IPs are the control list.",
              "alerts", "06  Alerts  ·  status, MITRE, analyst notes", "blocked", "07  Blocked IPs  ·  block and unblock")
    shot_pair("Models & reports", "Ch. 4  ·  Pages 08–09", "Measure the models. Export the proof.",
              "AI Models holds metrics and retrain. Reports: Detection, Incidents, Export, Compliance, Live Evidence.",
              "models", "08  AI Models  ·  accuracy, precision, F1", "reports", "09  Reports  ·  CSV, Live Evidence, compliance")
    shot_pair("Settings & incidents", "Ch. 4  ·  Pages 10 and 12", "Configure the platform. Case the alerts.",
              "Settings tabs: General, Telegram, AI, NDR, Clear, Company. Incidents group alerts into cases.",
              "settings", "10  Settings  ·  General, Telegram, AI, NDR, Clear, Company", "incidents", "12  Incidents  ·  cases, owner, status")
    shot_pair("Assets & hunting", "Ch. 4  ·  Pages 13–14", "Know the host. Hunt the trail.",
              "Assets: inventory, topology, UEBA. Hunting: query, PCAP, protocol.",
              "assets", "13  Assets  ·  inventory, topology, host profile", "hunting", "14  Hunting  ·  hunt, PCAP, DNS/TLS")
    shot_pair("SOC Ops", "Ch. 4  ·  Page 11", "Map the technique. Run the playbook.",
              "SOC Ops tabs: MITRE, SOAR, online learning, server health, ML registry.",
              "mitre", "11  SOC Ops  ·  MITRE ATT&CK", "soar", "11  SOC Ops  ·  SOAR playbooks")
    shot_pair("Explain & act", "Ch. 4  ·  Pages 15–16", "Ask why. Then allow, block, or wait.",
              "Response: IOC, allow/block, approvals, webhooks, sensors. Copilot explains with stored XAI.",
              "resp", "15  Response  ·  IOC, allow, block, sensors", "copilot", "16  Copilot  ·  natural-language assist")

    # Results
    s = add("Results", "Say the numbers, then the disclaimer in the same breath. Honesty is a grade.")
    d.kicker(s, "Ch. 5  ·  Testing and evaluation")
    d.h1(s, "Table 9 and Table 11. Same measured holdout.")
    d.accent_card(s, 0.4, 1.35, 3.05, 2.15, "click_r1")
    d.text(s, 0.6, 1.5, 2.7, 0.35, "RANDOM FOREST", 12, True, GOLD, name="with_r1a")
    d.text(s, 0.6, 1.9, 2.7, 0.7, "97.8%", 36, True, WHITE, name="with_r1b")
    d.text(s, 0.6, 2.7, 2.7, 0.5, "Accuracy & weighted F1", 13, False, MUTED, name="with_r1c")
    d.accent_card(s, 3.6, 1.35, 3.05, 2.15, "click_r2")
    d.text(s, 3.8, 1.5, 2.7, 0.35, "XGBOOST", 12, True, GOLD, name="with_r2a")
    d.text(s, 3.8, 1.9, 2.7, 0.7, "98.0%", 36, True, WHITE, name="with_r2b")
    d.text(s, 3.8, 2.7, 2.7, 0.5, "Accuracy & weighted F1", 13, False, MUTED, name="with_r2c")
    d.accent_card(s, 6.8, 1.35, 3.05, 2.15, "click_r3")
    d.text(s, 7.0, 1.5, 2.7, 0.35, "ISOLATION FOREST", 12, True, GOLD, name="with_r3a")
    d.text(s, 7.0, 1.9, 2.7, 0.7, "93.9%", 36, True, WHITE, name="with_r3b")
    d.text(s, 7.0, 2.7, 2.7, 0.5, "Normal vs Attack proxy", 13, False, MUTED, name="with_r3c")
    d.accent_card(s, 10.0, 1.35, 2.9, 2.15, "click_r4")
    d.text(s, 10.2, 1.5, 2.55, 0.35, "EVAL SET", 12, True, GOLD, name="with_r4a")
    d.text(s, 10.2, 1.9, 2.55, 0.7, "12k", 36, True, WHITE, name="with_r4b")
    d.text(s, 10.2, 2.7, 2.55, 0.5, "rows · 39 features · 6 classes", 12, False, MUTED, name="with_r4c")
    chart = EVAL / "fig_model_performance.png"
    if chart.exists():
        d.picture(s, chart, 0.45, 3.65, 7.4, "click_chart")
    d.rect(s, 8.05, 3.65, 4.85, 2.85, CARD, name="click_disc")
    d.text(s, 8.25, 3.8, 4.5, 0.35, "ACADEMIC SCOPE", 12, True, GOLD, name="with_da")
    d.text(s, 8.25, 4.25, 4.5, 2.0, "Same numbers as the paper Table 9 / Table 11. 12,000-row CICIDS-style holdout, 39 features, six classes, 75/25 split, random_state=42. Not an official CICIDS2017 leaderboard. AE/LSTM/Hybrid run in the engine.", 14, False, WHITE, name="with_db")

    # Strengths
    s = add("Strengths", "Own the strengths. Then the weaknesses. The jury respects both.")
    d.kicker(s, "Ch. 6  ·  Strengths")
    d.h1(s, "What this work already does well.")
    strengths = [
        ("End-to-end SOC", "16 real pages from capture to Copilot — not a mock dashboard."),
        ("Hybrid AI", "RF, XGBoost, Isolation Forest, Autoencoder, LSTM, plus online SGD."),
        ("Same pipeline", "Live traffic and simulated campaigns share detect → alert → respond."),
        ("Operations", "MITRE, SOAR, Telegram, stored XAI, Company trial, human-approval mode."),
        ("Honest science", "97.8% RF / 98.0% XGB on a 12k CICIDS-style holdout, with an explicit academic-scope disclaimer."),
        ("Defendable stack", "Local FastAPI, Streamlit, SQLite (Postgres optional) — launchable from the desktop."),
    ]
    for i, (t, b) in enumerate(strengths):
        x, y = 0.45 + (i % 3) * 4.2, 1.4 + (i // 3) * 2.55
        d.accent_card(s, x, y, 4.0, 2.35, f"click_sk{i}")
        d.text(s, x + 0.25, y + 0.25, 3.5, 0.45, t, 18, True, GOLD, name=f"with_sk{i}t")
        d.text(s, x + 0.25, y + 0.85, 3.5, 1.2, b, 15, False, WHITE, name=f"with_sk{i}b")

    # Contributions
    s = add("Contributions", "Four bullets only. This is what you added beyond a model file.")
    d.kicker(s, "Conclusion  ·  Contributions")
    d.h1(s, "What this project adds.")
    contrib = [
        ("Working NDR 2.0", "A local, versioned SOC product — not a single notebook."),
        ("Hybrid AI loop", "Five detectors plus online learning that preserves core models."),
        ("Operations layer", "MITRE, SOAR, incidents, hunting, Copilot, Telegram, and response control."),
        ("Honest evaluation", "12k CICIDS-style holdout, stored XAI, and live-versus-simulation evidence."),
    ]
    for i, (t, b) in enumerate(contrib):
        x = 0.45 + i * 3.2
        d.accent_card(s, x, 1.5, 3.05, 4.7, f"click_c{i}")
        d.text(s, x + 0.2, 1.8, 2.65, 0.5, f"0{i+1}", 26, True, GOLD, name=f"with_c{i}n")
        d.text(s, x + 0.2, 2.5, 2.65, 1.1, t, 20, True, WHITE, name=f"with_c{i}t")
        d.text(s, x + 0.2, 3.8, 2.65, 2.0, b, 15, False, MUTED, name=f"with_c{i}b")

    # Limitations
    s = add("Limitations", "State scope before they ask. It signals academic maturity, not a weak product.")
    d.kicker(s, "Ch. 6  ·  Limitations")
    d.h1(s, "Scope we state on purpose.")
    limits = [
        "Evaluation uses a 12,000-row CICIDS-style set — not an official CICIDS2017 leaderboard.",
        "Prototype is local (SQLite default, optional Postgres). It is not a clustered enterprise NDR.",
        "Live capture needs OS permissions and a selected NIC; results depend on that room.",
        "Company trial (roles, backup, SIEM hooks) is a laboratory/office pilot, not Active Directory SSO.",
        "Headline numbers focus on RF / XGBoost / Isolation Forest; AE and LSTM are in the engine.",
    ]
    for i, t in enumerate(limits):
        y = 1.35 + i * 0.95
        d.accent_card(s, 0.45, y, 12.4, 0.85, f"click_lim{i}")
        d.text(s, 0.75, y + 0.2, 11.9, 0.52, t, 15, False, WHITE, name=f"with_lim{i}")

    # Vision
    s = add("Vision", "Vision is future work with a direction. Keep it credible.")
    d.kicker(s, "Ch. 6  ·  Future work")
    d.h1(s, "Where this NDR should go next.")
    fut = [
        ("Public benchmarks", "Retrain and report on full CICIDS2017 and UNSW-NB15 with published splits."),
        ("Explain every hit", "Keep stored per-alert XAI; add SHAP plots for high-severity alerts."),
        ("Sensor fleet", "SPAN/TAP or multi-host sensors beyond the local NIC and 24/7 capture daemon."),
        ("Hardening", "Active Directory SSO and a full container stack — roles, backup, and SIEM hooks already started."),
    ]
    for i, (t, b) in enumerate(fut):
        x, y = 0.45 + (i % 2) * 6.4, 1.45 + (i // 2) * 2.5
        d.accent_card(s, x, y, 6.15, 2.25, f"click_fu{i}")
        d.text(s, x + 0.3, y + 0.3, 5.6, 0.5, t, 22, True, GOLD, name=f"with_fu{i}t")
        d.text(s, x + 0.3, y + 0.95, 5.6, 1.0, b, 16, False, WHITE, name=f"with_fu{i}b")

    # Demo path (paper 7.6 / Chapter 8)
    s = add("Demo path", "If live demo fails, this slide is the backup story. Memorize the seven steps.")
    d.kicker(s, "Ch. 7–8  ·  Demonstration")
    d.h1(s, "Seven clicks the jury should see.")
    steps = [
        "Launch run.bat / desktop shortcut",
        "Sign in as administrator",
        "Start live capture",
        "Launch a mixed campaign",
        "Run hybrid detection",
        "Open alert + stored XAI",
        "Block + Telegram + Copilot",
    ]
    for i, t in enumerate(steps):
        if i < 4:
            x, y = 0.45 + i * 3.2, 1.45
        else:
            x, y = 2.05 + (i - 4) * 3.2, 3.95
        d.rect(s, x, y, 3.05, 2.2, CARD, GOLD, name=f"click_dm{i}")
        d.text(s, x + 0.2, y + 0.3, 2.65, 0.55, f"{i+1:02d}", 28, True, GOLD, name=f"with_dm{i}n")
        d.text(s, x + 0.2, y + 1.0, 2.65, 0.9, t, 16, True, WHITE, name=f"with_dm{i}t")

    # Conclusion
    s = add("Conclusion", "Three sentences. Then stop. Invite questions.")
    d.kicker(s, "Conclusion")
    d.h1(s, "A working academic NDR. Measured. Scoped.")
    cons = [
        "The problem was detecting known and unknown attacks without a full capture–detect–respond loop.",
        "The solution is the implemented local NDR: hybrid AI, 16 SOC pages, stored XAI, Telegram, and optional company-trial hardening.",
        "The evidence is the running system plus Table 11: RF 97.8%, XGB 98.0%, IF 93.9% on a 12,000-row CICIDS-style holdout — not a CICIDS2017 leaderboard.",
    ]
    for i, t in enumerate(cons):
        y = 1.5 + i * 1.55
        d.rect(s, 0.45, y, 1.1, 1.3, GOLD, name=f"click_co{i}")
        d.text(s, 0.45, y + 0.4, 1.1, 0.5, f"{i+1}", 28, True, NAVY, align=PP_ALIGN.CENTER, name=f"with_co{i}n")
        d.accent_card(s, 1.75, y, 11.1, 1.3, f"with_co{i}c")
        d.text(s, 2.05, y + 0.35, 10.6, 0.7, t, 18, False, WHITE, name=f"with_co{i}t")

    # 23 References
    s = add("References", "Do not read all. Point to the paper. Mention MITRE, CICIDS, XGBoost, RF.")
    d.kicker(s, "Selected references")
    d.h1(s, "Sources the jury will recognise.")
    refs = [
        "Breiman, L. (2001) Random forests. Machine Learning.",
        "Chen, T. and Guestrin, C. (2016) XGBoost. KDD.",
        "Liu, F.T., Ting, K.M. and Zhou, Z.-H. (2008) Isolation Forest. ICDM.",
        "Sharafaldin et al. CICIDS2017 — Canadian Institute for Cybersecurity.",
        "MITRE ATT&CK® framework. MITRE Corporation.",
        "Scikit-learn, XGBoost, PyTorch, Streamlit, FastAPI documentation.",
        "Full Harvard-style list: AUCE research paper (desktop / appendix).",
    ]
    d.accent_card(s, 0.45, 1.4, 12.4, 5.15, "click_ref")
    d.multilines(s, 0.75, 1.6, 11.9, 4.8, [f"•  {r}" for r in refs], 16, WHITE, "with_ref", 10)

    # 24 Thank you
    s = add("Q&A", "Smile. Repeat the question. Answer in 30–40 seconds. Offer a live page if relevant.")
    if LOGO.exists():
        d.picture(s, LOGO, 5.95, 1.15, 1.45, "auto_logo2")
    d.text(s, 0.5, 2.8, 12.3, 0.9, "Thank you.", 54, True, WHITE, align=PP_ALIGN.CENTER, name="auto_ty")
    d.text(s, 0.5, 3.75, 12.3, 0.45, "Questions, challenges, and live walkthroughs are welcome.", 18, False, GOLD, align=PP_ALIGN.CENTER, name="click_ty1")
    d.text(s, 0.5, 4.4, 12.3, 0.4, "Mahmoud Talal Kanaan  ·  AUCE 2026  ·  AI Network Analyzer v1.0.0", 16, False, MUTED, align=PP_ALIGN.CENTER, name="click_ty2")
    d.text(s, 0.5, 5.1, 12.3, 0.35, "Desktop launcher  ·  Dashboard :8501  ·  API :8000", 14, False, CYAN, align=PP_ALIGN.CENTER, name="click_ty3")

    total = len(S)
    for i, sld in enumerate(S, 1):
        section = d.slides[i - 1][1]
        d.footer(sld, i, total, section)
        add_fade_transition(sld)
    apply_notes(d.prs, d.notes)

    wrote = []
    for dest in (OUT_PROJ, OUT_DESK):
        dest.parent.mkdir(parents=True, exist_ok=True)
        try:
            d.prs.save(str(dest))
            wrote.append(dest)
            print("WROTE", dest)
        except PermissionError:
            alt = dest.with_name(dest.stem + "_ALIGNED.pptx")
            d.prs.save(str(alt))
            wrote.append(alt)
            print("LOCKED", dest, "WROTE", alt)
    return wrote[0] if wrote else OUT_PROJ


# PowerPoint animation IDs
FX_FLY, FX_FADE, FX_WIPE = 2, 10, 22
FX_ZOOM, FX_FLOAT, FX_RISE = 23, 30, 34
FX_ASCEND, FX_FADED_ZOOM, FX_GLIDE, FX_EXPAND = 39, 48, 49, 50
DIR_LEFT, DIR_RIGHT, DIR_TOP, DIR_BOTTOM = 1, 2, 3, 4
AFTER, WITH = 3, 2

# Unique motion + transition per slide (1-based)
SLIDE_STYLE = {
    1:  {"card": FX_FLOAT, "title": FX_RISE, "pic": FX_FADED_ZOOM, "trans": 3845},
    2:  {"card": FX_ASCEND, "title": FX_RISE, "pic": FX_FADE, "trans": 3852},
    3:  {"card": FX_FLY, "title": FX_WIPE, "pic": FX_FADE, "trans": 3852, "dir": DIR_LEFT},
    4:  {"card": FX_WIPE, "title": FX_RISE, "pic": FX_FADE, "trans": 3853, "dir": DIR_LEFT},
    5:  {"card": FX_EXPAND, "title": FX_FADE, "pic": FX_FADE, "trans": 3845},
    6:  {"card": FX_FADED_ZOOM, "title": FX_RISE, "pic": FX_FADE, "trans": 3852},
    7:  {"card": FX_FLOAT, "title": FX_ASCEND, "pic": FX_FADE, "trans": 3852},
    8:  {"card": FX_GLIDE, "title": FX_WIPE, "pic": FX_GLIDE, "trans": 3845},
    9:  {"card": FX_FADED_ZOOM, "title": FX_RISE, "pic": FX_FADED_ZOOM, "trans": 3845},
    10: {"card": FX_EXPAND, "title": FX_FADE, "pic": FX_FADE, "trans": 3852},
    11: {"card": FX_ASCEND, "title": FX_RISE, "pic": FX_FADE, "trans": 3852},
    12: {"card": FX_FADED_ZOOM, "title": FX_FADE, "pic": FX_FADED_ZOOM, "trans": 3845},
    13: {"card": FX_GLIDE, "title": FX_WIPE, "pic": FX_FADED_ZOOM, "trans": 3845},
    14: {"card": FX_FLY, "title": FX_RISE, "pic": FX_FADED_ZOOM, "trans": 3845, "dir": DIR_LEFT},
    15: {"card": FX_FLY, "title": FX_WIPE, "pic": FX_FADED_ZOOM, "trans": 3845, "dir": DIR_RIGHT},
    16: {"card": FX_FLOAT, "title": FX_RISE, "pic": FX_FADED_ZOOM, "trans": 3845},
    17: {"card": FX_FADED_ZOOM, "title": FX_RISE, "pic": FX_FADED_ZOOM, "trans": 3852},
    18: {"card": FX_ASCEND, "title": FX_WIPE, "pic": FX_FADE, "trans": 3852},
    19: {"card": FX_FLOAT, "title": FX_RISE, "pic": FX_FADE, "trans": 3853},
    20: {"card": FX_WIPE, "title": FX_FADE, "pic": FX_FADE, "trans": 3852, "dir": DIR_LEFT},
    21: {"card": FX_EXPAND, "title": FX_ASCEND, "pic": FX_FADE, "trans": 3853},
    22: {"card": FX_RISE, "title": FX_RISE, "pic": FX_FADE, "trans": 3845},
    23: {"card": FX_FADE, "title": FX_WIPE, "pic": FX_FADE, "trans": 3845},
    24: {"card": FX_RISE, "title": FX_RISE, "pic": FX_FADED_ZOOM, "trans": 3845},
}


def _fx_for(shp, style: dict) -> int:
    n = str(shp.Name).lower()
    try:
        is_pic = int(shp.Type) == 13
    except Exception:
        is_pic = False
    if is_pic or "logo" in n:
        return style.get("pic", FX_FADED_ZOOM)
    if "title" in n or n.endswith("_ty") or n == "auto_ty":
        return style.get("title", FX_RISE)
    if "kicker" in n or "uni" in n:
        return FX_WIPE
    if n.startswith("click_") or n.startswith("auto_"):
        return style.get("card", FX_FLOAT)
    return FX_FADE


def _add_eff(seq, shp, effect, trigger, duration=1.35, delay=0.0, direction=None):
    eff = seq.AddEffect(shp, effect, 0, trigger)
    if direction is not None:
        try:
            eff.EffectParameters.Direction = direction
        except Exception:
            pass
    try:
        eff.Timing.Duration = duration
    except Exception:
        pass
    if delay:
        try:
            eff.Timing.TriggerDelayTime = delay
        except Exception:
            pass
    try:
        eff.Timing.SmoothEnd = True
    except Exception:
        pass
    try:
        eff.Timing.Duration = duration
    except Exception:
        pass
    return eff


def polish_with_powerpoint(path: Path) -> None:
    import win32com.client

    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    powerpoint.Visible = 1
    pres = powerpoint.Presentations.Open(str(path), WithWindow=True)

    for i in range(1, pres.Slides.Count + 1):
        sl = pres.Slides(i)
        style = SLIDE_STYLE.get(i) or list(SLIDE_STYLE.values())[(i - 1) % len(SLIDE_STYLE)]
        tr = sl.SlideShowTransition
        try:
            tr.EntryEffect = style["trans"]
        except Exception:
            tr.EntryEffect = 3845
        try:
            tr.Duration = 0.85
        except Exception:
            pass
        tr.AdvanceOnClick = True
        tr.AdvanceOnTime = False

        seq = sl.TimeLine.MainSequence
        try:
            while seq.Count > 0:
                seq(1).Delete()
        except Exception:
            pass

        named = []
        for si in range(1, sl.Shapes.Count + 1):
            shp = sl.Shapes(si)
            n = str(shp.Name)
            if n.startswith(("auto_", "click_", "with_")):
                named.append(shp)

        first = True
        for shp in named:
            n = str(shp.Name)
            effect = _fx_for(shp, style)
            direction = style.get("dir")
            if n.startswith("click_") and i in (3, 14) and "click_" in n:
                # alternate fly-in sides on problem / screenshot slides
                try:
                    idx = int("".join(ch for ch in n if ch.isdigit()) or "0")
                    direction = DIR_LEFT if idx % 2 == 0 else DIR_RIGHT
                except Exception:
                    direction = style.get("dir")
            if n.startswith("with_"):
                _add_eff(seq, shp, FX_FADE, WITH, duration=1.15, delay=0.12, direction=None)
            elif first:
                _add_eff(seq, shp, effect, AFTER, duration=1.35, delay=0.25, direction=direction)
                first = False
            else:
                # new beat — like a video cut, about 1s slower than before
                _add_eff(seq, shp, effect, AFTER, duration=1.25, delay=0.18, direction=direction)

        for j in range(1, seq.Count + 1):
            try:
                if float(seq(j).Timing.Duration) < 1.1:
                    seq(j).Timing.Duration = 1.3
            except Exception:
                continue

    pres.Save()
    print("POLISHED", path)
    powerpoint.Activate()


if __name__ == "__main__":
    out = build()
    try:
        polish_with_powerpoint(out)
    except Exception as exc:
        print("COM polish skipped:", type(exc).__name__, exc)

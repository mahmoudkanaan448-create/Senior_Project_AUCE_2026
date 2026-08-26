"""Align defense PPTX text exactly with RESEARCH_ONLY_FINAL_UPDATE.docx — text only."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = Path(r"c:\Users\mohamad\OneDrive\Desktop\AI_Network_Analyzer_Senior_Project_Defense.pptx")

REPLACES: dict[str, str] = {
    # Agenda / modules
    "Stack and all 16 SOC pages": "Fifteen dashboard modules (thesis)",
    "16 pages: alerts, incidents, hunting, reports": "Fifteen modules: capture, detect, alert, hunt, respond",
    "16 real pages from capture to Copilot — not a mock dashboard.": (
        "Fifteen dashboard modules from capture to Copilot — not a mock dashboard."
    ),
    "The solution is the implemented local NDR: hybrid AI, 16 SOC pages, stored XAI, Telegram, and optional company-trial hardening.": (
        "The solution is the implemented local NDR: hybrid AI, fifteen dashboard modules, stored XAI, Telegram, and optional company-trial hardening."
    ),
    # Tables (thesis numbering)
    "12k holdout — RF/XGB/IF only (Table 11)": "12k holdout — RF/XGB/IF only (Tables 11 & 21)",
    "Table 9 and Table 11. Same measured holdout.": "Table 11 and Table 21. Same measured holdout.",
    "Same numbers as Table 9 / Table 11. 12,000-row CICIDS-style holdout, 39 features, six classes, 75/25, random_state=42. Not a CICIDS2017 leaderboard. AE/LSTM/Hybrid: implemented; not scored on this holdout.": (
        "Same numbers as Table 11 / Table 21. 12,000-row CICIDS-style holdout, 39 features, six classes, 75/25, random_state=42. "
        "Not a CICIDS2017 leaderboard. AE/LSTM/Hybrid: implemented but not on holdout."
    ),
    # Objectives — match Table 2 exactly
    "Monitor network traffic and extract useful security features (39 CICIDS-style flow features).": (
        "Monitor network traffic and extract useful security features."
    ),
    "Detect known attacks using supervised models: Random Forest and XGBoost.": (
        "Detect known attacks using supervised ML models such as Random Forest and XGBoost."
    ),
    "Unknown anomalies: Isolation Forest (holdout) and Autoencoder (not on holdout).": (
        "Detect unknown anomalies using Isolation Forest and Autoencoder."
    ),
    "LSTM sequences (window_size=10, 10×39); implemented but not on holdout.": (
        "Use LSTM to analyze sequential behavior and support future attack prediction."
    ),
    "Support incident response: IP blocking, Telegram, webhooks, and reports (no email).": (
        "Support incident response using IP blocking, Telegram alerts, webhooks, and reports (no email channel)."
    ),
    # Hybrid / holdout wording
    "RF/XGB/IF measured; AE/LSTM implemented (not Table 11); Hybrid = fuse_decisions().": (
        "RF/XGB/IF measured (Table 21); AE/LSTM/Hybrid not on holdout; Hybrid = fuse_decisions()."
    ),
    "Headline metrics: RF/XGB/IF only. AE, LSTM, and Hybrid fusion are not Table 11 holdout scores.": (
        "Headline metrics: RF/XGB/IF only (Table 21). AE, LSTM, and Hybrid are not on holdout."
    ),
    "12,000-row holdout: RF 97.8% / XGB 98.0% / IF 93.9%. AE/LSTM/Hybrid not on holdout.": (
        "12,000-row holdout (Table 21): RF 97.8% / XGB 98.0% / IF 93.9%. AE/LSTM/Hybrid not on holdout."
    ),
    # Demo — Table 24
    "CH. 7  ·  USER MANUAL / DEMO": "CH. 7  ·  DEFENSE DEMO (Table 24)",
    "Live path from the thesis: capture → features → fusion → alert → XAI → MITRE → response.": (
        "Defense flow from Table 24: dataset → models → live demo path."
    ),
    "Launch run.bat": "Dataset: 12,000 rows · 39 features · stratified holdout",
    "Sign in (administrator)": "Models: RF/XGB/IF measured; AE/LSTM/Hybrid honest",
    "Capture / CSV → 39 features": "Traffic/CSV → Features → Preprocess",
    "Threat Simulation: DoS or mixed campaign": "Models → Hybrid fusion → Alert",
    "Models vote → fuse_decisions()": "XAI → MITRE → SOAR/Telegram",
    "Alert + stored XAI + MITRE": "Database → Dashboard (run.bat :8000 / :8501)",
    "Response / Telegram → dashboard proof": "Threat Simulation on TEST-NET (optional live step)",
}


def set_para_text(paragraph, new: str) -> None:
    runs = paragraph.runs
    if runs:
        runs[0].text = new
        for r in runs[1:]:
            r.text = ""
    else:
        paragraph.text = new


def walk(shape, fn) -> None:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            fn(p)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    fn(p)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for g in shape.shapes:
            walk(g, fn)


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    prs = Presentation(str(PPTX))
    n = 0

    def apply(p) -> None:
        nonlocal n
        cur = "".join(r.text or "" for r in p.runs) if p.runs else (p.text or "")
        if cur in REPLACES:
            set_para_text(p, REPLACES[cur])
            n += 1

    for slide in prs.slides:
        for sh in slide.shapes:
            walk(sh, apply)

    prs.save(str(PPTX))
    print("Updated paragraphs:", n)
    print("Saved:", PPTX)


if __name__ == "__main__":
    main()

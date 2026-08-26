"""Update defense PPTX text only — no layout, animations, or slide count changes."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPTX = Path(r"c:\Users\mohamad\OneDrive\Desktop\AI_Network_Analyzer_Senior_Project_Defense.pptx")

# Exact current paragraph text -> thesis SUPERVISOR_REVISED wording
REPLACES: dict[str, str] = {
    "CH. 7  ·  DEMONSTRATION": "CH. 7  ·  USER MANUAL / DEMO",
    "Seven clicks the jury should see.": "Live path from the thesis: capture → features → fusion → alert → XAI → MITRE → response.",
    "Launch run.bat / desktop shortcut": "Launch run.bat",
    "Sign in as administrator": "Sign in (administrator)",
    "Start live capture": "Capture / CSV → 39 features",
    "Launch a mixed campaign": "Threat Simulation: DoS or mixed campaign",
    "Run hybrid detection": "Models vote → fuse_decisions()",
    "Open alert + stored XAI": "Alert + stored XAI + MITRE",
    "Block + Telegram + Copilot": "Response / Telegram → dashboard proof",
    "Use LSTM to analyze sequential behavior and support attack prediction.": (
        "LSTM sequences (window_size=10, 10×39); implemented but not on holdout."
    ),
    "Detect unknown anomalies using Isolation Forest and Autoencoder.": (
        "Unknown anomalies: Isolation Forest (holdout) and Autoencoder (not on holdout)."
    ),
    "Live capture + hybrid models + threat score + severity.": (
        "Live capture + RF/XGB/IF (+ AE/LSTM in engine) → fuse_decisions() → threat score."
    ),
    "Reports, live-versus-simulation evidence, and the 12k holdout.": (
        "12,000-row holdout: RF 97.8% / XGB 98.0% / IF 93.9%. AE/LSTM/Hybrid not on holdout."
    ),
    "Accuracy table, retrain from CSV, online learning status": (
        "Holdout metrics (RF/XGB/IF), retrain from CSV, online SGD status"
    ),
    "08  AI Models  ·  accuracy, precision, F1": (
        "08  AI Models  ·  RF 97.8% / XGB 98.0% / IF 93.9%"
    ),
    "Questions, challenges, and live walkthroughs are welcome.": (
        "Q&A: why 12,000 not 720? why these 39 features? leakage? why high scores? why Hybrid fusion?"
    ),
}


def set_para_text(paragraph, new: str) -> None:
    runs = paragraph.runs
    if runs:
        runs[0].text = new
        for r in runs[1:]:
            r.text = ""
    else:
        paragraph.text = new


def walk(shape, fn) -> None:
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            fn(p)
    if shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                for p in cell.text_frame.paragraphs:
                    fn(p)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for g in shape.shapes:
            walk(g, fn)


def main() -> None:
    sys.stdout.reconfigure(encoding="utf-8")
    prs = Presentation(str(PPTX))
    n = 0

    def apply(p) -> None:
        nonlocal n
        cur = "".join(r.text or "" for r in p.runs) if p.runs else (p.text or "")
        if cur in REPLACES:
            set_para_text(p, REPLACES[cur])
            n += 1
            return

    for slide in prs.slides:
        for sh in slide.shapes:
            walk(sh, apply)

    prs.save(str(PPTX))
    print("Updated paragraphs:", n)
    print("Saved:", PPTX)


if __name__ == "__main__":
    main()
